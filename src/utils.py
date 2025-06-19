import PyPDF2
from pptx import Presentation

def read_file(file_path):
        if file_path.endswith(".pdf"):
            with open (file_path, "rb") as file:
                try:
                    pdf_reader = PyPDF2.PdfReader(file)
                    text = ""
                    for page in pdf_reader.pages:
                        text += page.extract_text()
                    return text
                except Exception as e:
                    raise Exception("error reading the PDF file")
        
        elif file_path.endswith(".txt"):
            with open(file_path, "r", encoding='utf-8') as file:
                text = file.read()
            return text
    
        elif file_path.endswith(".pptx"):
            prs = Presentation(file_path)
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return text
    
        else:
            raise Exception("unsupported file format, only PDF, text, and PowerPoint files are supported")

def split_text_into_chunks(text, max_length=1500):
    sentences = text.split('.')
    chunks = []
    current_chunk_sentences = [] 
    current_chunk_length = 0

    for sentence in sentences:
        sentence = sentence.strip() 

        if not sentence: 
            continue

        sentence_with_period = sentence + '.'
        sentence_length = len(sentence_with_period)

        if current_chunk_length + sentence_length <= max_length:
            current_chunk_sentences.append(sentence_with_period)
            current_chunk_length += sentence_length
        else:
            if current_chunk_sentences: 
                chunks.append("".join(current_chunk_sentences).strip()) 
            current_chunk_sentences = [sentence_with_period] 
            current_chunk_length = sentence_length

    if current_chunk_sentences:
        chunks.append("".join(current_chunk_sentences).strip())

    return chunks


def split_list_into_chunks(input_list, n):
  chunked_list = []
  for i in range(0, len(input_list), n):
    chunked_list.append(input_list[i:i+n])
  return chunked_list
